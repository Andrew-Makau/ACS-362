#!/usr/bin/env python3
"""
Generate a Word report and PowerPoint deck with embedded cash-flow charts.

Features:
- Driver-based 24-month forecast with Base/Upside/Downside scenarios
- Embedded line charts (Base, Scenarios) and WC components chart
- Optional PowerPoint template (slide master/theme) support
- CLI arguments to override defaults (start month, horizon, capex, etc.)
- Headless-safe plotting (matplotlib Agg), basic logging, and error handling

Usage (local):
  python scripts/generate_docs.py --start 2025-11 --horizon 24 \
    --opening-cash 2000000 --capex "3:500000,12:1200000" \
    --template "assets/ACS362A_Template.pptx"

If no CLI arguments are supplied, sensible defaults are used.
"""

from __future__ import annotations

import argparse
import logging
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from datetime import datetime

import numpy as np
import pandas as pd

import matplotlib
# Headless backend (safe for CI and servers)
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from docx import Document
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.util import Inches as PPTInches

# ============== Defaults (can be overridden via CLI) ==============
DEFAULT_COURSE = "ACS 362A: Database Systems Group Project"
DEFAULT_INSTRUCTOR = "Japheth Mursi, Ph.D."
DEFAULT_TEAM: List[str] = [
    "22-2207 – Makau Andrew Nzyoki",
    "22-2203 – Wanjihia Serena Wairati",
    "22-2245 – Dorothy Vugutsa",
    "22-2337 – Lewis Kagia",
    "23-2189 – Kitavi Christoph",
]
DEFAULT_CONTEXT = "Nairobi-based consumer electronics retailer (store + online)"
DEFAULT_CURRENCY = "KES"
DEFAULT_HORIZON_MONTHS = 24
DEFAULT_OPENING_CASH = 2_000_000
DEFAULT_CAPEX_SCHEDULE = {3: 500_000, 12: 1_200_000}
DEFAULT_START_REV = 5_000_000
DEFAULT_COGS_PCT = 0.40
DEFAULT_FIXED_OPEX = 1_800_000
DEFAULT_VAR_OPEX_PCT = 0.08
DEFAULT_DEPR_LIFE_MONTHS = 60
DEFAULT_TAX_RATE = 0.30
DEFAULT_SCENARIOS = {
    "Base":     {"growth": 0.015, "dso": 45, "dpo": 30, "inv_days": 60},
    "Upside":   {"growth": 0.030, "dso": 35, "dpo": 35, "inv_days": 50},
    "Downside": {"growth": 0.005, "dso": 55, "dpo": 25, "inv_days": 70},
}


@dataclass
class Paths:
    root: Path
    deliverables: Path
    charts: Path
    data: Path
    report: Path
    deck: Path
    csv: Path
    template: Optional[Path] = None


@dataclass
class Config:
    course: str = DEFAULT_COURSE
    instructor: str = DEFAULT_INSTRUCTOR
    team: List[str] = None  # set in __post_init__
    context: str = DEFAULT_CONTEXT
    currency: str = DEFAULT_CURRENCY
    horizon_months: int = DEFAULT_HORIZON_MONTHS
    opening_cash: float = DEFAULT_OPENING_CASH
    capex_schedule: Dict[int, float] = None  # set in __post_init__
    start_month: datetime = None            # set in __post_init__
    start_rev: float = DEFAULT_START_REV
    cogs_pct: float = DEFAULT_COGS_PCT
    fixed_opex: float = DEFAULT_FIXED_OPEX
    var_opex_pct: float = DEFAULT_VAR_OPEX_PCT
    depr_life_months: int = DEFAULT_DEPR_LIFE_MONTHS
    tax_rate: float = DEFAULT_TAX_RATE
    scenarios: Dict[str, Dict[str, float]] = None  # set in __post_init__

    def __post_init__(self):
        if self.team is None:
            self.team = list(DEFAULT_TEAM)
        if self.capex_schedule is None:
            self.capex_schedule = dict(DEFAULT_CAPEX_SCHEDULE)
        if self.start_month is None:
            today = datetime.today()
            self.start_month = datetime(today.year, today.month, 1)
        if self.scenarios is None:
            self.scenarios = dict(DEFAULT_SCENARIOS)


def make_paths(template: Optional[str] = None) -> Paths:
    root = Path(__file__).resolve().parents[1]
    deliverables = root / "deliverables"
    charts = deliverables / "charts"
    data = deliverables / "data"
    deliverables.mkdir(parents=True, exist_ok=True)
    charts.mkdir(parents=True, exist_ok=True)
    data.mkdir(parents=True, exist_ok=True)

    p = Paths(
        root=root,
        deliverables=deliverables,
        charts=charts,
        data=data,
        report=deliverables / "Cash_Flow_Forecasting_Report_Group2.docx",
        deck=deliverables / "Cash_Flow_Forecasting_Presentation_Group2.pptx",
        csv=data / "cash_series.csv",
    )
    if template:
        t = (root / template) if not template.startswith(("~", "/", "\\")) else Path(template).expanduser()
        p.template = t if t.exists() else None
        if template and p.template is None:
            logging.warning("Template not found at %s; proceeding without template.", t)
    return p


def parse_capex(text: str) -> Dict[int, float]:
    """
    Parse CapEx schedule from a string like "3:500000,12:1200000".
    """
    schedule: Dict[int, float] = {}
    if not text:
        return schedule
    for part in text.split(","):
        part = part.strip()
        if not part:
            continue
        if ":" not in part:
            raise ValueError(f"Invalid capex item '{part}'. Expected format 'month:amount'.")
        m_s, amt_s = part.split(":", 1)
        month = int(m_s.strip())
        amount = float(amt_s.strip().replace("_", "").replace(",", ""))
        if month < 1:
            raise ValueError(f"Invalid month {month} in capex schedule.")
        schedule[month] = amount
    return schedule


def build_schedule(cfg: Config, scenario: Dict[str, float]) -> pd.DataFrame:
    months = pd.date_range(cfg.start_month, periods=cfg.horizon_months, freq="MS")
    df = pd.DataFrame({"month": months})

    # Revenue
    rev = [cfg.start_rev]
    for _ in range(1, cfg.horizon_months):
        rev.append(rev[-1] * (1 + scenario["growth"]))
    df["revenue"] = rev

    # COGS, Gross Margin
    df["cogs"] = df["revenue"] * cfg.cogs_pct
    df["gross_profit"] = df["revenue"] - df["cogs"]

    # Opex
    df["opex"] = cfg.fixed_opex + cfg.var_opex_pct * df["revenue"]
    df["ebitda"] = df["gross_profit"] - df["opex"]

    # CapEx
    df["capex"] = 0.0
    for m, amt in cfg.capex_schedule.items():
        if 1 <= m <= cfg.horizon_months:
            df.loc[m - 1, "capex"] = float(amt)

    # Depreciation: straight-line for each CapEx over DEPR_LIFE_MONTHS
    dep = np.zeros(cfg.horizon_months, dtype=float)
    for m, amt in cfg.capex_schedule.items():
        start_idx = m - 1
        life = min(cfg.depr_life_months, cfg.horizon_months - start_idx)
        if life > 0:
            dep[start_idx:start_idx + life] += float(amt) / cfg.depr_life_months
    df["depreciation"] = dep
    df["ebit"] = df["ebitda"] - df["depreciation"]

    # Taxes (only on positive EBIT for simplicity)
    df["tax"] = df["ebit"].clip(lower=0) * cfg.tax_rate
    df["net_income"] = df["ebit"] - df["tax"]

    # Working capital (monthly approximation)
    dso = scenario["dso"]
    dpo = scenario["dpo"]
    inv_days = scenario["inv_days"]
    df["ar"] = df["revenue"] * (dso / 30.0)
    df["ap"] = df["cogs"] * (dpo / 30.0)
    df["inventory"] = df["cogs"] * (inv_days / 30.0)
    df["delta_ar"] = df["ar"].diff().fillna(df["ar"])
    df["delta_ap"] = df["ap"].diff().fillna(df["ap"])
    df["delta_inv"] = df["inventory"].diff().fillna(df["inventory"])
    df["delta_wc"] = df["delta_ar"] + df["delta_inv"] - df["delta_ap"]

    # Cash flow
    df["cfo"] = df["net_income"] + df["depreciation"] - df["delta_wc"]
    df["net_cash_flow"] = df["cfo"] - df["capex"]

    # Cash balance
    cash = [cfg.opening_cash + df.loc[0, "net_cash_flow"]]
    for i in range(1, cfg.horizon_months):
        cash.append(cash[-1] + df.loc[i, "net_cash_flow"])
    df["closing_cash"] = cash

    return df


def build_all(cfg: Config) -> Dict[str, pd.DataFrame]:
    return {name: build_schedule(cfg, sc) for name, sc in cfg.scenarios.items()}


def save_data(paths: Paths, dfs: Dict[str, pd.DataFrame]) -> pd.DataFrame:
    out = pd.DataFrame({"month": dfs["Base"]["month"]})
    for name in ["Base", "Upside", "Downside"]:
        out[f"{name}_closing_cash"] = dfs[name]["closing_cash"].round(2)
    out["Base_net_cash_flow"] = dfs["Base"]["net_cash_flow"].round(2)
    out.to_csv(paths.csv, index=False)
    logging.info("Wrote %s", paths.csv)
    return out


def _style_plot(title: str, y_label: str):
    plt.title(title)
    plt.xlabel("Month")
    plt.ylabel(y_label)
    plt.grid(True, alpha=0.3)
    plt.tight_layout()


def chart_closing_cash_base(paths: Paths, cfg: Config, df: pd.DataFrame) -> Path:
    plt.figure(figsize=(10, 5))
    plt.plot(df["month"], df["closing_cash"], label="Closing Cash (Base)", color="#1f77b4", linewidth=2)
    _style_plot("Closing Cash — Base Case", f"Cash ({cfg.currency})")
    p = paths.charts / "closing_cash_base.png"
    plt.savefig(p, dpi=200)
    plt.close()
    logging.info("Wrote %s", p)
    return p


def chart_scenarios(paths: Paths, cfg: Config, dfs: Dict[str, pd.DataFrame]) -> Path:
    plt.figure(figsize=(10, 5))
    for name, color in zip(["Base", "Upside", "Downside"], ["#1f77b4", "#2ca02c", "#d62728"]):
        plt.plot(dfs[name]["month"], dfs[name]["closing_cash"], label=name, linewidth=2, color=color)
    plt.legend()
    _style_plot("Closing Cash — Base vs Upside vs Downside", f"Cash ({cfg.currency})")
    p = paths.charts / "closing_cash_scenarios.png"
    plt.savefig(p, dpi=200)
    plt.close()
    logging.info("Wrote %s", p)
    return p


def chart_wc_components(paths: Paths, cfg: Config, df: pd.DataFrame, month_idx: int = 11) -> Path:
    # ΔAR, ΔInventory, −ΔAP for Month 12 (index 11) by default
    m = min(max(month_idx, 0), len(df) - 1)
    components = {
        "ΔAR": df.loc[m, "delta_ar"],
        "ΔInventory": df.loc[m, "delta_inv"],
        "−ΔAP": -df.loc[m, "delta_ap"],
    }
    plt.figure(figsize=(8, 5))
    bars = plt.bar(list(components.keys()), list(components.values()), color=["#9467bd", "#8c564b", "#ff7f0e"])
    for b in bars:
        y = b.get_height()
        plt.text(b.get_x() + b.get_width()/2, y, f"{y:,.0f}", ha="center", va="bottom", fontsize=9)
    _style_plot(f"Working Capital Components — Month {m+1}", f"Cash impact ({cfg.currency})")
    p = paths.charts / "wc_components.png"
    plt.savefig(p, dpi=200)
    plt.close()
    logging.info("Wrote %s", p)
    return p


def build_report(paths: Paths, cfg: Config, chart_paths: Dict[str, Path]) -> None:
    doc = Document()
    # Title
    doc.add_heading("Cash Flow Forecasting Model — Group 2", level=0)
    doc.add_paragraph(f"Course: {cfg.course}")
    doc.add_paragraph(f"Instructor: {cfg.instructor}")
    doc.add_paragraph("Prepared by: Group 2 — " + "; ".join(cfg.team))
    doc.add_paragraph(f"Date: {datetime.today():%B %d, %Y} • Version: v1.0")

    doc.add_heading("Executive Summary", level=1)
    doc.add_paragraph(
        "Objective: Build a transparent monthly cash flow forecasting model to project cash position, "
        "identify funding needs, and evaluate scenarios."
    )
    doc.add_paragraph(
        "Scope: 24‑month forecast in KES with Base / Upside / Downside scenarios; working capital (AR/AP/Inventory), "
        "CapEx/Depreciation, and optional debt."
    )
    doc.add_paragraph(f"Case Context: {cfg.context}")

    doc.add_heading("Key Charts", level=1)
    doc.add_paragraph("Closing cash — Base case:")
    doc.add_picture(str(chart_paths["base_line"]), width=Inches(6))
    doc.add_paragraph("Closing cash — Base vs Upside vs Downside:")
    doc.add_picture(str(chart_paths["scenarios_line"]), width=Inches(6))
    doc.add_paragraph("Working capital components (ΔAR, ΔInventory, −ΔAP):")
    doc.add_picture(str(chart_paths["wc_components"]), width=Inches(6))

    doc.add_heading("Methodology (summary)", level=1)
    doc.add_paragraph(
        "Revenue via monthly growth; COGS as % of revenue; Opex with fixed and variable components; "
        "working capital using DSO/DPO/InventoryDays; CapEx and straight‑line depreciation; simple taxes; "
        "cash flow = Net Income + Depreciation − ΔWC − CapEx."
    )

    doc.add_heading("Scenario Assumptions", level=1)
    for name, cfg_sc in cfg.scenarios.items():
        doc.add_paragraph(
            f"{name}: growth {cfg_sc['growth']*100:.1f}%/mo, DSO {cfg_sc['dso']}d, DPO {cfg_sc['dpo']}d, "
            f"Inventory {cfg_sc['inv_days']}d"
        )

    doc.add_heading("Notes & Recommendations", level=1)
    doc.add_paragraph(
        "Focus on reducing DSO, negotiating DPO, and phasing CapEx to smooth the cash trough. "
        "Update actuals monthly and refresh scenarios quarterly."
    )

    doc.save(str(paths.report))
    logging.info("Wrote %s", paths.report)


def _get_layout(prs: Presentation, preferred_index: int, fallback_index: int = 1):
    try:
        return prs.slide_layouts[preferred_index]
    except Exception:
        return prs.slide_layouts[fallback_index]


def _first_body_placeholder(slide) -> Optional[object]:
    # Find first non-title placeholder for content
    for ph in slide.placeholders:
        try:
            ph_type = ph.placeholder_format.type
        except Exception:
            continue
        # Title type is 1, skip it
        if ph_type != 1:
            return ph
    return None


def build_deck(paths: Paths, cfg: Config, chart_paths: Dict[str, Path]) -> None:
    # Load template if provided, else default
    prs = Presentation(str(paths.template)) if paths.template else Presentation()

    # Title slide
    title_layout = _get_layout(prs, 0, 0)
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "Cash Flow Forecasting Model (Group 2)"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = f"{cfg.course}\nInstructor: {cfg.instructor}\nTeam: " + "; ".join(cfg.team)

    def add_slide(title: str, bullets: List[str], notes: Optional[str] = None, layout_idx: int = 1):
        layout = _get_layout(prs, layout_idx, 1)
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = title
        body = _first_body_placeholder(s) or s.shapes.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = b
        if notes:
            s.notes_slide.notes_text_frame.text = notes

    add_slide(
        "Agenda",
        ["Objective & Scope", "Approach & Assumptions", "Model Overview", "Results & Scenarios", "Sensitivities, Risks", "Recommendations"],
        "Keep to ~10–12 minutes; Q&A at the end.",
    )
    add_slide(
        "Objective & Scope",
        ["24‑month cash forecast in KES", "Revenue, margins, opex, working capital, CapEx", "Base / Upside / Downside scenarios"],
        "Why 24 months: shows seasonality; concise for class.",
    )
    add_slide(
        "Approach (Model Architecture)",
        ["Driver‑based; centralized Assumptions", "Schedules: Revenue, Opex, Working Capital, CapEx/Dep, Debt", "Consolidated Cash Flow; scenario switcher"],
        "No hardcoded numbers in formulas.",
    )
    add_slide(
        "Key Assumptions (Base)",
        ["Growth: 1.5% monthly", "COGS: 40% of revenue", "DSO/DPO/InvDays: 45 / 30 / 60", "Tax: 30%; Interest: 8% if debt"],
        None,
    )

    # Chart slides — try "Title Only" layout (often index 5) else fallback
    title_only_layout = 5
    for slide_title, key in [
        ("Results: Base Closing Cash", "base_line"),
        ("Scenario Comparison: Closing Cash", "scenarios_line"),
        ("Working Capital Components", "wc_components"),
    ]:
        layout = _get_layout(prs, title_only_layout, 0)
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = slide_title
        s.shapes.add_picture(str(chart_paths[key]), PPTInches(1), PPTInches(1.5), height=PPTInches(4.5))

    add_slide(
        "Recommendations",
        ["Reduce DSO; negotiate DPO", "Phase CapEx after trough", "Maintain buffer via revolving facility", "Update actuals monthly; refresh scenarios quarterly"],
        None,
    )

    prs.save(str(paths.deck))
    logging.info("Wrote %s", paths.deck)


def run(cfg: Config, paths: Paths) -> None:
    dfs = build_all(cfg)
    _ = save_data(paths, dfs)

    # Charts
    base_line = chart_closing_cash_base(paths, cfg, dfs["Base"])
    scenarios_line = chart_scenarios(paths, cfg, dfs)
    wc_components = chart_wc_components(paths, cfg, dfs["Base"], month_idx=11)

    chart_paths = {"base_line": base_line, "scenarios_line": scenarios_line, "wc_components": wc_components}

    # Docs
    build_report(paths, cfg, chart_paths)
    build_deck(paths, cfg, chart_paths)

    print(
        "Generated:\n"
        f"- {paths.report}\n- {paths.deck}\n- {paths.csv}\n"
        f"Charts in {paths.charts}"
    )


def _parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate Cash Flow Report and Deck with Charts.")
    parser.add_argument("--start", type=str, default=None, help="Start month YYYY-MM (default = current month).")
    parser.add_argument("--horizon", type=int, default=DEFAULT_HORIZON_MONTHS, help="Forecast horizon in months.")
    parser.add_argument("--opening-cash", type=float, default=DEFAULT_OPENING_CASH, help="Opening cash amount.")
    parser.add_argument("--capex", type=str, default="", help='CapEx schedule, e.g., "3:500000,12:1200000".')
    parser.add_argument("--course", type=str, default=DEFAULT_COURSE)
    parser.add_argument("--instructor", type=str, default=DEFAULT_INSTRUCTOR)
    parser.add_argument("--context", type=str, default=DEFAULT_CONTEXT)
    parser.add_argument("--currency", type=str, default=DEFAULT_CURRENCY)
    parser.add_argument("--template", type=str, default=None, help="Relative or absolute path to PPTX template.")
    parser.add_argument("--start-rev", type=float, default=DEFAULT_START_REV, help="Starting monthly revenue (month 1).")
    parser.add_argument("--log-level", type=str, default="INFO", help="Logging level (DEBUG, INFO, WARNING, ERROR).")
    return parser.parse_args()


def _build_config_from_args(args: argparse.Namespace) -> Tuple[Config, Paths]:
    logging.basicConfig(level=getattr(logging, args.log_level.upper(), logging.INFO), format="%(levelname)s: %(message)s")
    # Start month
    if args.start:
        try:
            start_month = datetime.strptime(args.start + "-01", "%Y-%m-%d")
        except ValueError as e:
            raise SystemExit(f"Invalid --start '{args.start}'. Expected YYYY-MM.") from e
    else:
        today = datetime.today()
        start_month = datetime(today.year, today.month, 1)

    # CapEx schedule
    schedule = parse_capex(args.capex) if args.capex else dict(DEFAULT_CAPEX_SCHEDULE)

    cfg = Config(
        course=args.course,
        instructor=args.instructor,
        team=list(DEFAULT_TEAM),
        context=args.context,
        currency=args.currency,
        horizon_months=args.horizon,
        opening_cash=args.opening_cash,
        capex_schedule=schedule,
        start_month=start_month,
        start_rev=args.start_rev,
        cogs_pct=DEFAULT_COGS_PCT,
        fixed_opex=DEFAULT_FIXED_OPEX,
        var_opex_pct=DEFAULT_VAR_OPEX_PCT,
        depr_life_months=DEFAULT_DEPR_LIFE_MONTHS,
        tax_rate=DEFAULT_TAX_RATE,
        scenarios=dict(DEFAULT_SCENARIOS),
    )
    paths = make_paths(template=args.template)
    return cfg, paths


def main() -> None:
    args = _parse_args()
    cfg, paths = _build_config_from_args(args)
    try:
        run(cfg, paths)
    except Exception as exc:
        logging.exception("Generation failed: %s", exc)
        raise SystemExit(1) from exc


if __name__ == "__main__":
    main()